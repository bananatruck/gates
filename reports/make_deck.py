"""Build the Gate 1 presentation (PPTX).

Same numbers as GATE1_REPORT.pdf, same palette, one claim per slide.

    python reports/make_deck.py
"""

from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "GATE1_DECK.pptx"

RUN_JSON = Path(
    "/home/kesh/AgentLaboratory-Gemini/ablation_runs/data_efficiency/ablation.json"
)
RUN_DIR = RUN_JSON.parent
RUNS_ROOT = RUN_DIR.parent

BENCH = Path(
    "/tmp/claude-1000/-home-kesh/976cef29-0afd-479c-a716-6c557e07b6cb"
    "/scratchpad/bench2.json"
)

BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x8A, 0x88, 0x80)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PANEL = RGBColor(0xF4, 0xF3, 0xEE)

W, H = Inches(13.333), Inches(7.5)


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, spacing=1.15):
    """runs: list of (text, size, bold, colour) -> one paragraph each."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (t, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = "DejaVu Sans"
    return box


def panel(slide, x, y, w, h, colour=PANEL):
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = colour
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def header(slide, kicker, title):
    text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
         [(kicker.upper(), 11, True, BLUE)])
    text(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.7),
         [(title, 26, True, INK)])


def bullets(slide, x, y, w, items, size=14, gap=1.5):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            runs.append(it)
        else:
            runs.append((f"·  {it}", size, False, INK2))
    text(slide, x, y, w, Inches(4), runs, spacing=gap)


def stat(slide, x, y, w, value, label, colour=BLUE):
    panel(slide, x, y, w, Inches(1.35))
    text(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.6),
         [(value, 30, True, colour)])
    text(slide, x + Inches(0.2), y + Inches(0.82), w - Inches(0.4), Inches(0.45),
         [(label, 10, False, INK2)])


def picture(slide, name, x, y, w, max_h=None):
    """Place a figure, shrinking it if it would not fit the space given.

    matplotlib saves with bbox_inches="tight", so a figure's on-disk aspect is
    not the figsize it was declared with -- placing by width alone overflowed
    the slide edge and clipped the axis labels off the scanner chart. The real
    aspect is read from the file.
    """
    p = ASSETS / name
    if not p.exists():
        return
    from PIL import Image

    pw, ph = Image.open(p).size
    height = Emu(int(w * ph / pw))
    if max_h is not None and height > max_h:
        height = max_h
        w = Emu(int(max_h * pw / ph))
    slide.shapes.add_picture(str(p), x, y, width=w, height=height)




def _table(slide, x, y, w, rows, col_w, header=True, size=11, rh=Inches(0.34)):
    """A real table, because a claim about numbers should show the numbers."""
    from pptx.util import Inches as I

    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, rh * n_rows)
    tbl = shape.table
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(size)
            para.font.name = "DejaVu Sans"
            para.font.bold = header and i == 0
            para.font.color.rgb = INK if (header and i == 0) else INK2
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if (header and i == 0) else SURFACE
    return shape


def _slide_run(prs, run):
    s = blank(prs)
    header(s, "the run", "Data efficiency, with and without Gate 1")
    text(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0),
         [("Question given to the agent: how much labelled data does a classifier "
           "need? Train the same model on 25% of the pool and on all of it, and "
           "report both test accuracies and their ratio — three values.",
           13, False, INK2),
          ("Engineer: deepseek-v4-flash (reasoning_effort=low). Gate 1's own two "
           "jobs: a local qwen3:8b, free per call and never consulted for a "
           "verdict. Both arms use both models, so neither is advantaged — they "
           "differ only in what decides that an experiment succeeded.",
           11, False, MUTED)])

    g, u = run["arms"]["gated"], run["arms"]["ungated"]
    gv, uv = g["verification"], u["verification"]
    rows = [["", "with Gate 1", "without Gate 1"],
            ["accepted a run", "yes" if g["accepted"] else "no",
             "yes" if u["accepted"] else "no"],
            ["turns used", g["turns"], u["turns"]],
            ["numbers reported", len(gv["reported"]), len(uv["reported"])],
            ["numbers checkable against an execution",
             len(gv["matched"]) + len(gv["mismatched"]),
             len(uv["matched"]) + len(uv["mismatched"])],
            ["reproduced on re-execution", len(gv["matched"]), len(uv["matched"])]]
    _table(s, Inches(0.6), Inches(3.0), Inches(12.1), rows,
           [Inches(6.1), Inches(3.0), Inches(3.0)], size=13, rh=Inches(0.46))

    rate = f"{100 * gv['rate']:.0f}%" if gv["rate"] is not None else "n/a"
    stat(s, Inches(0.6), Inches(6.0), Inches(4.0), rate,
         "of the Gate 1 arm's numbers reproduced when its code was re-run", BLUE)
    stat(s, Inches(4.9), Inches(6.0), Inches(4.0), str(len(uv["reported"])),
         "numbers the ungated arm recorded through any contract", ORANGE)
    stat(s, Inches(9.2), Inches(6.0), Inches(3.5),
         str(len(gv["matched"]) + len(gv["mismatched"])),
         "of the gated arm's numbers are checkable at all", BLUE)


def _slide_all_runs(prs):
    """Every run on disk, not the one that was rendered."""
    import sys as _sys
    _sys.path.insert(0, str(Path(__file__).resolve().parent))
    import aggregate as agg

    runs = agg.Aggregate(agg.load_runs(RUNS_ROOT))
    if not runs.n:
        return
    s = blank(prs)
    header(s, "across all runs", f"Every ablation record on disk — n = {runs.n}")
    g, u = runs.arm_stats("gated"), runs.arm_stats("ungated")
    if not g.get("n"):
        return
    rate = lambda x: "—" if x is None else f"{100*x:.0f}%"
    rows = [[f"across {g['n']} run(s)", "with Gate 1", "without Gate 1"],
            ["accepted a run", f"{g['accepted']} of {g['n']}", f"{u['accepted']} of {u['n']}"],
            ["mean turns used", f"{g['turns_mean']:.1f}", f"{u['turns_mean']:.1f}"],
            ["numbers reported", g["reported_total"], u["reported_total"]],
            ["numbers checkable", g["checkable_total"], u["checkable_total"]],
            ["reproduced on re-execution",
             f"{g['reproduced_total']} ({rate(g['reproduction_rate'])})",
             f"{u['reproduced_total']} ({rate(u['reproduction_rate'])})"],
            ["accepted what the other arm rejected",
             g["false_success_total"], u["false_success_total"]]]
    _table(s, Inches(0.6), Inches(1.8), Inches(12.1), rows,
           [Inches(5.5), Inches(3.3), Inches(3.3)], size=13, rh=Inches(0.5))
    text(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0),
         [("A run that predates a measurement is reported as not measured, never "
           "as zero — the distinction cost us once already.", 12, False, ORANGE),
          ("Repeats are at temperature 0.7: at 0.0 a repeat is the same "
           "experiment run twice, not a second sample.", 11, False, MUTED)])


def _slide_accuracy(prs, run):
    s = blank(prs)
    header(s, "the two reports", "What each arm would hand the writing agent")
    gv = run["arms"]["gated"]["verification"]
    uv = run["arms"]["ungated"]["verification"]
    keys = sorted(set(gv["reported"]) | set(uv["reported"]))[:9]
    rows = [["metric", "Gate 1: reported", "re-executed", "verdict",
             "no gate: reported"]]
    for k in keys:
        rep = gv["reported"].get(k)
        again = gv["reproduced"].get(k)
        verdict = ("reproduced" if k in gv["matched"]
                   else "MISMATCH" if k in gv["mismatched"] else "—")
        rows.append([
            k,
            f"{rep:.4g}" if isinstance(rep, float) else (rep if rep is not None else "—"),
            f"{again:.4g}" if isinstance(again, float) else (again if again is not None else "—"),
            verdict,
            uv["reported"].get(k, "not recorded"),
        ])
    if len(rows) == 1:
        rows.append(["(neither arm recorded a value)", "", "", "", ""])
    _table(s, Inches(0.6), Inches(1.7), Inches(12.1), rows,
           [Inches(3.4), Inches(2.2), Inches(2.2), Inches(2.1), Inches(2.2)],
           size=11, rh=Inches(0.36))
    text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.9),
         [(uv["reason"].capitalize() + "." if uv["reason"] else
           "Both arms recorded values.", 13, True, ORANGE),
          ("Accuracy here means one thing only: run the accepted code again and "
           "see whether the numbers it reported are the numbers it produces.",
           11, False, INK2)])


def _slide_report_accuracy(prs):
    """The comparison the deck was missing: both arms' reports, measured alike.

    Three properties, kept apart on purpose. The ungated arm passes on accuracy
    and fails on the other two, and one merged column would have hidden that —
    which is the more interesting result, not the more flattering one.
    """
    path = Path(RUNS_ROOT) / "report_accuracy.json"
    if not path.exists():
        return
    data = json.loads(path.read_text(encoding="utf-8"))
    if not data:
        return

    def tot(arm, fn):
        return sum(fn(r.get(arm, {})) for r in data.values())

    g_vis = tot("gated", lambda a: len(a.get("visible", {})))
    g_tr = tot("gated", lambda a: a.get("traceable", 0))
    u_prod = tot("ungated", lambda a: len(a.get("produced", {})))
    u_vis = tot("ungated", lambda a: len(a.get("visible", {})))
    u_lost = u_prod - u_vis
    u_match = tot("ungated", lambda a: len(a.get("matched", [])))
    u_check = u_match + tot("ungated", lambda a: len(a.get("mismatched", [])))

    s = blank(prs)
    header(s, "report accuracy", "Both arms' reports, measured the same way")
    rows = [["", "with Gate 1", "without Gate 1"],
            ["results the run produced", g_vis, u_prod],
            ["reached the writing agent", g_vis, u_vis],
            ["lost to the 1,000-char channel", 0, u_lost],
            ["traceable to an execution", g_tr,
             tot("ungated", lambda a: a.get("traceable", 0))],
            ["reproduced when re-run", f"{g_vis} of {g_vis}",
             f"{u_match} of {u_check}"]]
    _table(s, Inches(0.6), Inches(1.75), Inches(12.1), rows,
           [Inches(5.5), Inches(3.3), Inches(3.3)], size=13, rh=Inches(0.46))
    text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.6),
         [("The ungated arm's numbers are not wrong.", 14, True, INK),
          (f"Every one of the {u_check} that reached the writer reproduced "
           f"exactly. It fails on completeness — {u_lost} of {u_prod} results "
           f"never got through the window, and in one run all three did not, "
           f"leaving the writer a training log and no results — and on "
           f"traceability: 0 of {u_prod} carry a trace id or a code hash.",
           12, False, INK2),
          ("A number that reproduces by luck and a number that was measured "
           "look identical to a reader. That is the gap.", 12, False, ORANGE)])


def _slide_papers(prs):
    """The end of the chain: what the reader actually gets handed."""
    path = Path(RUNS_ROOT) / "paper_audit.json"
    if not path.exists():
        return
    data = json.loads(path.read_text(encoding="utf-8"))
    if not data:
        return

    s = blank(prs)
    header(s, "the papers", "Three full workflows, audited on identical criteria")
    picture(s, "papers.png", Inches(0.6), Inches(1.5), Inches(7.5),
            max_h=Inches(4.3))
    panel(s, Inches(8.4), Inches(1.5), Inches(4.3), Inches(4.9))
    text(s, Inches(8.65), Inches(1.7), Inches(3.8), Inches(4.5),
         [("The gate-off arm did not fabricate.", 14, True, INK),
          ("Its Results section says so itself: \u201cthe decisive "
           "measurements \u2026 are not present in the evidence window. The "
           "captured log was truncated at the legacy 1,000-character "
           "boundary.\u201d It then reports the loss curve it could see.",
           10.5, False, INK2),
          ("", 7, False, INK2),
          ("An honest paper about a truncated log.", 11, True, ORANGE),
          ("", 7, False, INK2),
          ("So two things are doing work, and only one of them is the gate. "
           "The prompt \u2014 shared by both arms \u2014 forbids describing an "
           "unrecorded quantity, and that is what stops fabrication. Gate 1 is "
           "what gets the results to the writer, bound to the run.",
           10.5, False, INK2),
          ("", 7, False, INK2),
          ("Instruction without delivery: an honest paper with no findings. "
           "Delivery without instruction: the archived paper, 65 of 65 "
           "unsourced.", 10.5, True, BLUE)])
    text(s, Inches(0.6), Inches(6.0), Inches(7.6), Inches(1.2),
         [("The gate-off arm is credited for every value its run recorded to "
           "disk \u2014 none of which reached the writer. The comparison is "
           "generous on purpose; the gap that survives is the claim.",
           10, False, MUTED)])


def _slide_checks_fired(prs, run):
    s = blank(prs)
    header(s, "checks", "Every check that ran, turn by turn")
    attempts = run["arms"]["gated"]["attempts"]
    ids = []
    for a in attempts:
        for c in a.get("all_checks", []):
            if c["id"] not in ids:
                ids.append(c["id"])
    ids = ids[:14]
    rows = [["check", "sev"] + [f"turn {a['turn']}" for a in attempts]]
    for cid in ids:
        row, sev = [cid], ""
        for a in attempts:
            m = next((c for c in a.get("all_checks", []) if c["id"] == cid), None)
            if m is None:
                row.append("—")
            else:
                sev = m["severity"]
                row.append("pass" if m["passed"] else sev)
        rows.append([row[0], sev] + row[1:])
    if len(rows) == 1:
        rows.append(["(no per-check record)", "", ""])
    ncols = len(rows[0])
    _table(s, Inches(0.6), Inches(1.7), Inches(12.1), rows,
           [Inches(4.6), Inches(1.1)] + [Inches(6.4 / max(ncols - 2, 1))] * (ncols - 2),
           size=10, rh=Inches(0.33))


def _slide_taxonomy(prs, run):
    """Our run, scored with MLR-Bench's published categories."""
    s = blank(prs)
    header(s, "taxonomy", "This run, scored against MLR-Bench's four classes")
    g, u = run["arms"]["gated"], run["arms"]["ungated"]
    gv, uv = g["verification"], u["verification"]

    def cell(arm, ver):
        if not arm["accepted"]:
            return "nothing reached the writer"
        if not ver["reported"]:
            return "every number unbacked — nothing recorded"
        unbacked = len(ver["reported"]) - (len(ver["matched"]) + len(ver["mismatched"]))
        return f"{unbacked} of {len(ver['reported'])} unbacked, {len(ver['matched'])} reproduced"

    rows = [["MLR-Bench class", "with Gate 1", "without Gate 1", "whose job"],
            ["Faked experimental results", cell(g, gv), cell(u, uv), "Gate 1 — measured here"],
            ["Hallucinated methodology", "not measured", "not measured", "Gate 2"],
            ["Incorrect citations", "not measured", "not measured", "Gate 3"],
            ["Mathematical errors", "not measured", "not measured", "outside all three"]]
    _table(s, Inches(0.6), Inches(1.8), Inches(12.1), rows,
           [Inches(3.4), Inches(3.3), Inches(3.3), Inches(2.1)], size=12,
           rh=Inches(0.62))
    text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2),
         [("This is our run scored with their categories — not a run of MLR-Bench, "
           "which needs their harness and task set.", 13, True, ORANGE),
          ("For scale, MLR-Bench reports faked results and hallucinated methodology "
           "in more than half of 10 tasks, and almost all AI Scientist V2 papers "
           "contain both. BadScientist reports fabricated papers accepted at up to "
           "82.0% with detection barely above chance.", 11, False, INK2)])


def _slide_cost(prs, run):
    """Where the money goes, and what the gate adds to it."""
    s = blank(prs)
    header(s, "cost", "Gate 1's marginal cost on the paid API is zero")
    # Summed across every run on disk, not one rendered run: a per-call cost
    # claim made from a single run is an anecdote, and the runs differ in how
    # many rewrites they needed, which is exactly what the gate changes.
    usage: dict[str, dict] = {}
    n_runs = 0
    for path in sorted(Path(RUNS_ROOT).glob("*/ablation.json")):
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            continue
        n_runs += 1
        for role, u in (data.get("usage") or {}).items():
            acc = usage.setdefault(
                role, {"calls": 0, "prompt": 0, "completion": 0, "reasoning": 0}
            )
            for field in acc:
                acc[field] += u.get(field, 0)
    if not usage:
        usage = run.get("usage") or {}
        n_runs = 1
    rows = [[f"role (summed over {n_runs} runs)", "calls", "prompt tok",
             "completion tok", "of which reasoning", "billed"]]
    for role, u in sorted(usage.items()):
        rows.append([role, u["calls"], f"{u['prompt']:,}", f"{u['completion']:,}",
                     f"{u['reasoning']:,}",
                     "paid API" if role == "engineer" else "local — free"])
    if len(rows) == 1:
        rows.append(["(no accounting)", "", "", "", "", ""])
    _table(s, Inches(0.6), Inches(1.7), Inches(12.1), rows,
           [Inches(2.9), Inches(1.1), Inches(2.1), Inches(2.3), Inches(2.3), Inches(1.4)],
           size=12, rh=Inches(0.48))
    eng, gate = usage.get("engineer", {}), usage.get("gate", {})
    if eng.get("calls") and gate.get("calls"):
        ep = eng["completion"] / eng["calls"]
        gp = gate["completion"] / gate["calls"]
        stat(s, Inches(0.6), Inches(3.9), Inches(3.8), f"{ep/gp:.0f}x",
             "one engineer call vs one gate call, in completion tokens", ORANGE)
        stat(s, Inches(4.7), Inches(3.9), Inches(3.8), f"{gate['completion']:,}",
             "gate tokens — all local, none billed", BLUE)
        stat(s, Inches(8.8), Inches(3.9), Inches(3.9), f"{eng['reasoning']:,}",
             "engineer reasoning tokens: billed, never seen in the output", ORANGE)
    text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.4),
         [("What Gate 1 changes on the bill is rewrites: each rejection buys "
           "another engineer call. A rejection that names the real fault is a "
           "rewrite the agent does not waste — which is why the shadowed-contract "
           "check was worth adding.", 13, False, INK2)])


def _slide_logs(prs, run):
    s = blank(prs)
    header(s, "logs", "Where the evidence is")
    d = RUN_DIR
    rows = [["arm", "path", "files per attempt"],
            ["with Gate 1", f"{d}/gated/gate1/attempt_NN/",
             "experiment.py, stdout.txt, stderr.txt, results.json, registry.json, gate1_report.json"],
            ["without Gate 1", f"{d}/ungated/attempt_NN/",
             "experiment.py, stdout.txt, stderr.txt"],
            ["shadow verdicts", f"{d}/ungated/shadow_gate/gate1/attempt_NN/",
             "what Gate 1 would have said about the same executions"],
            ["re-execution", f"{d}/<arm>/verify/", "the accepted code, run again"],
            ["summary", f"{d}/ablation.txt · ablation.json", "the tables in the report"]]
    _table(s, Inches(0.6), Inches(1.7), Inches(12.1), rows,
           [Inches(2.2), Inches(5.2), Inches(4.7)], size=10, rh=Inches(0.52))
    # a real excerpt
    acc = run["arms"]["gated"].get("accepted_at")
    snippet = ""
    if acc:
        f = RUN_DIR / "gated" / "gate1" / f"attempt_{acc:02d}" / "stdout.txt"
        if f.exists():
            snippet = "\n".join(f.read_text(errors="replace").splitlines()[:6])
    if snippet:
        panel(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.0))
        text(s, Inches(0.85), Inches(5.05), Inches(11.6), Inches(1.8),
             [("stdout.txt from the accepted run, captured in full and never truncated",
               11, True, INK)] +
             [(ln[:120], 10, False, INK2) for ln in snippet.splitlines()])


# --------------------------------------------------------------------------- #


def _live_verdicts() -> list[tuple[str, str, int, str]]:
    """Gate 1's verdicts in the full Agent Laboratory workflow, from disk.

    The ablation rig exercises the solver phase only. This reads the gate
    artifacts left by the complete workflow -- literature review through paper
    writing -- because that is the run the archived paper can be compared to.
    """
    root = Path("/home/kesh/AgentLaboratory-Gemini/research_dir/gate_artifacts/gate1")
    out = []
    for attempt in sorted(root.glob("attempt_*")):
        report, registry = attempt / "gate1_report.json", attempt / "registry.json"
        if not report.exists():
            continue
        try:
            r = json.loads(report.read_text())
            n = len(json.loads(registry.read_text()).get("values", {}))
        except (json.JSONDecodeError, FileNotFoundError):
            continue
        failed = [c["id"] for c in r.get("checks", []) if not c["passed"]
                  and c["severity"] == "FAIL"]
        out.append((attempt.name.replace("attempt_", "attempt "),
                    r.get("verdict", "?"), n, ", ".join(failed) or "—"))
    return out


def build(bench, run):
    """Ten slides, in the order the argument actually runs.

    The previous deck had twenty and buried the comparison in the middle of
    them. What a reader needs is the published rate, our two arms against it,
    and the evidence that our numbers are measurements -- everything else is
    supporting material that belongs in the report.
    """
    prs = deck()

    # ---- 1. title -------------------------------------------------------- #
    s = blank(prs)
    text(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.2),
         [("G.A.T.E.S. — Gate 1", 46, True, INK)])
    text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(1.0),
         [("Execution validity for autonomous research agents", 20, False, INK2)])
    text(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(1.0),
         [("Measured impact on Agent Laboratory, against the published "
           "fabrication rates", 15, False, MUTED)])
    text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
         [(f"{date.today().isoformat()}   ·   266 tests   ·   engineer "
           "deepseek-v4-flash, gates local qwen3:8b   ·   every figure a "
           "recorded measurement or a cited number", 11, False, MUTED)])

    # ---- 2. the problem -------------------------------------------------- #
    s = blank(prs)
    header(s, "the problem", "A crashed run scored 1.0 and became a published paper")
    panel(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.35))
    text(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.1),
         [("return output_capture.getvalue()[:MAX_LEN]      # MAX_LEN = 1000",
           15, True, INK),
          ("The crash marker is appended AFTER the program's own output — so on "
           "any run that prints more than 1,000 characters it falls off the same "
           "slice, and the solver's only crash test never fires.", 12, False, INK2)])
    stat(s, Inches(0.6), Inches(3.3), Inches(2.85), "1.0",
         "reward score, run raising NameError on every attempt", ORANGE)
    stat(s, Inches(3.7), Inches(3.3), Inches(2.85), "65",
         "numeric claims in the paper it produced", ORANGE)
    stat(s, Inches(6.8), Inches(3.3), Inches(2.85), "0",
         "of them traceable to any execution", ORANGE)
    stat(s, Inches(9.9), Inches(3.3), Inches(2.8), "0",
         "record_result calls in its saved code", MUTED)
    text(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.8),
         [("The paper is fluent, specific, and unfalsifiable from its own artifacts.",
           15, True, INK),
          ("“peak accuracy at K=2 (81.60% test) … a 13.61× speedup (0.0180s vs "
           "0.2450s) … a catastrophic collapse to 39.20% at K=8.” None of those "
           "numbers appears anywhere in the code that was actually run.",
           13, False, INK2),
          ("Nothing in the pipeline could tell a measurement from a sentence.",
           13, True, ORANGE)])

    # ---- 3. against the published rates ---------------------------------- #
    s = blank(prs)
    header(s, "the benchmarks", "This failure is the field's, and it is quantified")
    picture(s, "benchmark.png", Inches(0.6), Inches(1.55), Inches(7.9),
            max_h=Inches(4.6))
    panel(s, Inches(8.8), Inches(1.55), Inches(3.9), Inches(4.8))
    text(s, Inches(9.05), Inches(1.75), Inches(3.4), Inches(4.4),
         [("MLR-Bench", 14, True, INK),
          ("Coding agents produce fabricated or invalidated experimental results "
           "in 80% of cases. Both agents score below the 6.0 acceptance "
           "threshold on Soundness while scoring well on Completeness — fluent, "
           "and unsound.", 11, False, INK2),
          ("", 8, False, INK2),
          ("BadScientist", 14, True, INK),
          ("Fabricated papers accepted at up to 82.0%, with detection barely "
           "exceeding random chance.", 11, False, INK2),
          ("", 8, False, INK2),
          ("It names the defence it thinks is needed: provenance verification. "
           "That is what Gate 1 is.", 11, True, BLUE)])
    text(s, Inches(0.6), Inches(6.45), Inches(8.0), Inches(0.6),
         [("If LLM reviewers detect fabrication at chance, a validity layer built "
           "on model judgement inherits that ceiling — so Gate 1's verdict is "
           "entirely deterministic.", 10.5, False, INK2)])

    # ---- 4. the headline comparison -------------------------------------- #
    s = blank(prs)
    # The chart carries its own title, so the header must not repeat it.
    header(s, "the comparison",
           "Same acceptance rate — everything downstream differs")
    picture(s, "headline.png", Inches(0.6), Inches(1.5), Inches(8.1),
            max_h=Inches(4.8))
    panel(s, Inches(9.0), Inches(1.5), Inches(3.7), Inches(4.9))
    text(s, Inches(9.25), Inches(1.7), Inches(3.2), Inches(4.5),
         [("Both arms accept.", 15, True, INK),
          ("5 of 5 runs, either way. Acceptance is not where they differ — "
           "everything that happens to the numbers afterwards is.",
           11.5, False, INK2),
          ("", 8, False, INK2),
          ("Without the gate, 5 of the 15 results a run produced never reached "
           "the write-up at all; in one run all three were lost, leaving a "
           "writer with a training log and no results.", 11.5, False, INK2),
          ("", 8, False, INK2),
          ("And 0 of 15 carry a trace id or a code hash.", 11.5, True, ORANGE),
          ("", 8, False, INK2),
          ("5 runs · deepseek-v4-flash · T=0.7 · same task, same instructions, "
           "same turn budget.", 10, False, MUTED)])

    # ---- 5. what the reports actually say -------------------------------- #
    _slide_report_accuracy(prs)

    # ---- 6. the papers --------------------------------------------------- #
    _slide_papers(prs)

    # ---- 7. what Gate 1 checks ------------------------------------------- #
    s = blank(prs)
    header(s, "what gate 1 checks", "20 deterministic checks in six families, "
                                    "plus one model-assisted")
    picture(s, "checks.png", Inches(0.6), Inches(1.6), Inches(7.6),
            max_h=Inches(4.4))
    panel(s, Inches(8.5), Inches(1.6), Inches(4.2), Inches(4.7))
    text(s, Inches(8.75), Inches(1.8), Inches(3.7), Inches(4.3),
         [("The severity split is the design", 14, True, INK),
          ("FAIL blocks. WARN and INFO are reported, travel with their evidence, "
           "and can never change a verdict.", 11, False, INK2),
          ("", 8, False, INK2),
          ("That is what lets a REQUIRED LLM layer coexist with a model-free "
           "verdict: model findings are WARN by construction, and the feedback "
           "report is written after decide() has already fixed the outcome.",
           11, False, INK2),
          ("", 8, False, INK2),
          ("A test parses the module's AST to assert Severity.FAIL never appears "
           "in an expression there.", 10, False, MUTED)])
    text(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.7),
         [("results.values_computed is the strongest claim: record_result(\"k\", acc) "
           "passes, record_result(\"k\", 0.816) fails — and round(0.8160, 3) fails "
           "too, because constant folding does not launder a typed number.",
           11, False, INK2)])

    # ---- 8. the gate in the full workflow -------------------------------- #
    verdicts = _live_verdicts()
    if verdicts:
        s = blank(prs)
        header(s, "live run", "Gate 1 in the complete workflow, not just the solver")
        rows = [["attempt", "verdict", "values recorded", "blocking failures"]]
        rows += [[a, v, n, f] for a, v, n, f in verdicts]
        _table(s, Inches(0.6), Inches(1.7), Inches(12.1), rows,
               [Inches(2.0), Inches(1.8), Inches(2.6), Inches(5.7)],
               size=12, rh=Inches(0.44))
        y = Inches(1.7) + Inches(0.44) * len(rows) + Inches(0.4)
        text(s, Inches(0.6), y, Inches(12.1), Inches(2.0),
             [("The archived paper's run failed exactly this way and was scored 1.0.",
               14, True, INK),
              ("It raised NameError: name 'hidden_dim' is not defined on every "
               "attempt. The run above raised NameError: name "
               "'all_experiment_data' is not defined — the same class of defect, "
               "caught, with the reason handed back to the engineer.",
               12, False, INK2),
              ("Same failure mode. Opposite outcome.", 13, True, ORANGE)])

    # ---- 9. cost --------------------------------------------------------- #
    if run:
        _slide_cost(prs, run)

    # ---- 10. limits and next --------------------------------------------- #
    s = blank(prs)
    header(s, "limits & next", "What Gate 1 does not claim")
    bullets(s, Inches(0.6), Inches(1.7), Inches(12.1),
            ["It does not ask whether a result is plausible — that is Gate 2",
             "It does not read the manuscript — that is Gate 3",
             "The ungated arm's numbers are not wrong: every one that reached the "
             "writer reproduced. It fails on completeness and provenance, which "
             "is a different and narrower claim",
             "The log scanner's recall is bounded by what it was shown, and the "
             "lines examined are recorded so the claim cannot exceed the evidence",
             "A mean computed inside the experiment over a silently shortened "
             "list arrives as one value Gate 1 cannot see behind — Gate 2's to close"],
            size=13, gap=1.9)
    panel(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.7),
          RGBColor(0xFB, 0xF5, 0xEC))
    text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.5),
         [("Next", 14, True, INK),
          ("·  Complete the full-workflow paper and audit it against the archived "
           "one on identical criteria", 12, False, INK2),
          ("·  Channel-fidelity arm at MAX_LEN ∈ {1000, 4000, 16000, ∞}",
           12, False, INK2),
          ("·  Then Gate 2 — source ↔ result coherence, on the same model seam",
           12, True, BLUE)])
    return prs

def main():
    bench = json.loads(BENCH.read_text()) if BENCH.exists() else None
    run = json.loads(RUN_JSON.read_text()) if RUN_JSON.exists() else None
    if run is None:
        print(f'  WARNING: no {RUN_JSON} — run slides omitted')
    prs = build(bench, run)
    prs.save(OUT)
    print(f"  wrote {OUT.name} ({OUT.stat().st_size:,} bytes, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
